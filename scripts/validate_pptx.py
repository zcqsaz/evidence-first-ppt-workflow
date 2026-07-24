from __future__ import annotations

import argparse
import csv
import re
import sys
from collections import Counter, defaultdict
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

try:
    from .common import build_report, issue, load_json, resource_path, sha256_bytes, sha256_file, validate_json, write_report
except ImportError:
    from common import build_report, issue, load_json, resource_path, sha256_bytes, sha256_file, validate_json, write_report

IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".gif", ".webp", ".svg", ".emf", ".wmf"}
DEFAULT_BANNED_PHRASES = [
    "阅读要点",
    "课堂解读",
    "课堂观察",
    "理解要点",
    "案例定位",
    "课程示意",
    "图示所支持",
    "制作说明",
    "仅供 AI",
    "模块一",
    "模块二",
    "模块三",
]
FONT_ALIAS_GROUPS = [
    {"microsoft yahei", "microsoft yahei ui", "微软雅黑", "微软雅黑 ui"},
]


def default_config() -> dict[str, Any]:
    return {
        "schema_version": "1.0.0",
        "workflow_version": "1.0.0",
        "project": {"id": "unspecified", "name": "", "type": "other", "language": "zh-CN"},
        "scope": {"allowed_files": [], "forbidden_files": []},
        "presentation": {
            "expected_slide_count": None,
            "aspect_ratio": "16:9",
            "allowed_fonts": ["Microsoft YaHei", "微软雅黑"],
            "cover_slides": [1],
            "directory_slides": [2],
            "transition_slides": [],
            "page_number": {"required": False, "start_slide": 2, "center_tolerance_inches": 0.08, "bottom_zone_fraction": 0.35},
            "banned_phrases": DEFAULT_BANNED_PHRASES,
        },
        "qa": {
            "fail_on_out_of_bounds": True,
            "fail_on_text_overlap": True,
            "fail_on_nonzero_picture_crop": True,
            "fail_on_duplicate_evidence_media": True,
            "fail_on_banned_phrase": True,
            "fail_on_explicit_disallowed_font": True,
            "fail_on_unmapped_evidence_media": False,
            "fail_on_low_density": False,
            "ignore_template_media_repeated_on_fraction": 0.8,
            "text_overlap_smaller_area_fraction": 0.08,
            "bounds_tolerance_inches": 0.01,
            "crop_tolerance": 0.000001,
            "minimum_body_characters_without_visual": 180,
            "minimum_body_characters_with_visual": 80,
            "minimum_substantive_units": 4,
            "low_density_exempt_slides": [1, 2],
            "crop_exempt_slides": [],
            "overlap_exempt_slides": [],
            "out_of_bounds_exempt_slides": [],
            "duplicate_media_exempt_sha256": [],
        },
    }


def _deep_merge(base: dict[str, Any], override: dict[str, Any]) -> dict[str, Any]:
    merged = dict(base)
    for key, value in override.items():
        if isinstance(value, dict) and isinstance(merged.get(key), dict):
            merged[key] = _deep_merge(merged[key], value)
        else:
            merged[key] = value
    return merged


def _severity(config: dict[str, Any], flag: str) -> str:
    return "error" if config["qa"].get(flag, False) else "warning"


def _expanded_allowed_fonts(fonts: Iterable[str]) -> set[str]:
    allowed = {font.strip().casefold() for font in fonts if font.strip()}
    for aliases in FONT_ALIAS_GROUPS:
        normalized = {alias.casefold() for alias in aliases}
        if allowed & normalized:
            allowed.update(normalized)
    return allowed


def _shape_name(shape: Any) -> str:
    return getattr(shape, "name", "") or f"shape-{getattr(shape, 'shape_id', '?')}"


def _walk_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    for shape in shapes:
        yield shape
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(shape.shapes)


def _iter_text_frames(shape: Any) -> Iterable[Any]:
    if getattr(shape, "has_text_frame", False):
        yield shape.text_frame
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                yield cell.text_frame


def _shape_text(shape: Any) -> str:
    texts: list[str] = []
    for text_frame in _iter_text_frames(shape):
        texts.append(text_frame.text or "")
    return "\n".join(texts).strip()


def _slide_text(slide: Any) -> str:
    return "\n".join(filter(None, (_shape_text(shape) for shape in _walk_shapes(slide.shapes))))


def _iter_runs(shape: Any) -> Iterable[Any]:
    for text_frame in _iter_text_frames(shape):
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                yield run


def _run_typefaces(run: Any) -> set[str]:
    names: set[str] = set()
    if run.font.name:
        names.add(run.font.name)
    try:
        rpr = run._r.get_or_add_rPr()
        namespace = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        for child_name in ("latin", "ea", "cs", "sym"):
            child = rpr.find(namespace + child_name)
            if child is not None and child.get("typeface"):
                names.add(child.get("typeface"))
    except (AttributeError, TypeError):
        pass
    return {name.strip() for name in names if name and not name.startswith("+")}


def _rect(shape: Any) -> tuple[int, int, int, int]:
    left, top = int(shape.left), int(shape.top)
    return left, top, left + int(shape.width), top + int(shape.height)


def _overlap_fraction(a: Any, b: Any) -> float:
    ax1, ay1, ax2, ay2 = _rect(a)
    bx1, by1, bx2, by2 = _rect(b)
    width = max(0, min(ax2, bx2) - max(ax1, bx1))
    height = max(0, min(ay2, by2) - max(ay1, by1))
    intersection = width * height
    smaller = min(max(1, (ax2 - ax1) * (ay2 - ay1)), max(1, (bx2 - bx1) * (by2 - by1)))
    return intersection / smaller


def _top_level_text_shapes(slide: Any) -> list[Any]:
    result = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            continue
        if _shape_text(shape):
            result.append(shape)
    return result


def _picture_shapes(slide: Any) -> Iterable[Any]:
    for shape in _walk_shapes(slide.shapes):
        if getattr(shape, "shape_type", None) in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE}:
            yield shape


def _body_character_count(text: str) -> int:
    compact = re.sub(r"\s+", "", text)
    return len(compact)


def _substantive_unit_count(slide: Any) -> int:
    units = 0
    for shape in _walk_shapes(slide.shapes):
        for text_frame in _iter_text_frames(shape):
            for paragraph in text_frame.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue
                units += max(1, len([part for part in re.split(r"[。！？；;]", text) if len(part.strip()) >= 6]))
    return units


def _visual_count(slide: Any) -> int:
    count = 0
    for shape in _walk_shapes(slide.shapes):
        if getattr(shape, "shape_type", None) in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE, MSO_SHAPE_TYPE.CHART}:
            count += 1
        elif getattr(shape, "has_table", False):
            count += 1
    return count


def _build_asset_hash_index(asset_root: Path | None) -> tuple[dict[str, list[str]], dict[str, str]]:
    by_hash: defaultdict[str, list[str]] = defaultdict(list)
    metadata_hashes: dict[str, str] = {}
    if not asset_root or not asset_root.is_dir():
        return dict(by_hash), metadata_hashes
    for path in asset_root.rglob("*"):
        if path.is_file() and path.suffix.lower() in IMAGE_SUFFIXES:
            try:
                by_hash[sha256_file(path)].append(path.relative_to(asset_root).as_posix())
            except OSError:
                continue
    metadata_path = asset_root / "source_metadata.csv"
    if metadata_path.is_file():
        try:
            with metadata_path.open("r", encoding="utf-8-sig", newline="") as handle:
                for row in csv.DictReader(handle):
                    filename = (row.get("filename") or "").strip().replace("\\", "/")
                    if not filename:
                        continue
                    path = (asset_root / filename).resolve()
                    try:
                        path.relative_to(asset_root.resolve())
                    except ValueError:
                        continue
                    if path.is_file():
                        metadata_hashes[sha256_file(path)] = filename
        except (OSError, UnicodeError, csv.Error):
            pass
    return dict(by_hash), metadata_hashes


def validate_pptx(pptx_path: Path, config_path: Path | None = None, asset_root: Path | None = None) -> dict[str, Any]:
    pptx_path = pptx_path.resolve()
    issues: list[dict[str, Any]] = []
    config = default_config()
    if config_path:
        try:
            supplied = load_json(config_path.resolve())
        except (OSError, ValueError) as exc:
            return build_report("validate_pptx", pptx_path, [issue("error", "CONFIG_READ_ERROR", str(exc))])
        schema_errors = validate_json(supplied, resource_path("schemas/project_config.schema.json"))
        for message in schema_errors:
            issues.append(issue("error", "CONFIG_SCHEMA_VIOLATION", message))
        config = _deep_merge(config, supplied)

    if not pptx_path.is_file():
        return build_report("validate_pptx", pptx_path, issues + [issue("error", "PPTX_NOT_FOUND", f"PPTX not found: {pptx_path}")])
    try:
        presentation = Presentation(str(pptx_path))
    except Exception as exc:
        return build_report("validate_pptx", pptx_path, issues + [issue("error", "PPTX_OPEN_ERROR", f"PowerPoint package could not be opened: {exc}")])

    slide_count = len(presentation.slides)
    slide_width, slide_height = int(presentation.slide_width), int(presentation.slide_height)
    expected_count = config["presentation"].get("expected_slide_count")
    if expected_count is not None and slide_count != expected_count:
        issues.append(issue("error", "SLIDE_COUNT_MISMATCH", f"Expected {expected_count} slides, found {slide_count}."))

    expected_ratio = config["presentation"].get("aspect_ratio")
    actual_ratio = slide_width / slide_height
    target_ratios = {"16:9": 16 / 9, "4:3": 4 / 3}
    if expected_ratio in target_ratios and abs(actual_ratio - target_ratios[expected_ratio]) > 0.02:
        issues.append(issue("error", "ASPECT_RATIO_MISMATCH", f"Configured {expected_ratio}; actual ratio is {actual_ratio:.4f}."))

    qa = config["qa"]
    presentation_config = config["presentation"]
    allowed_fonts = _expanded_allowed_fonts(presentation_config.get("allowed_fonts", []))
    banned_phrases = presentation_config.get("banned_phrases", [])
    crop_exempt = set(qa.get("crop_exempt_slides", []))
    overlap_exempt = set(qa.get("overlap_exempt_slides", []))
    bounds_exempt = set(qa.get("out_of_bounds_exempt_slides", []))
    low_density_exempt = set(qa.get("low_density_exempt_slides", []))
    bounds_tolerance = int(Inches(float(qa.get("bounds_tolerance_inches", 0.01))))
    crop_tolerance = float(qa.get("crop_tolerance", 0.000001))
    overlap_threshold = float(qa.get("text_overlap_smaller_area_fraction", 0.08))
    crop_count = 0
    bounds_count = 0
    overlap_count = 0
    banned_count = 0
    disallowed_font_count = 0
    low_density_count = 0
    slide_metrics: list[dict[str, Any]] = []
    media_usage: defaultdict[str, list[dict[str, Any]]] = defaultdict(list)

    for slide_number, slide in enumerate(presentation.slides, start=1):
        text = _slide_text(slide)
        body_characters = _body_character_count(text)
        substantive_units = _substantive_unit_count(slide)
        visual_count = _visual_count(slide)
        slide_metrics.append({
            "slide": slide_number,
            "characters": body_characters,
            "substantive_units": substantive_units,
            "visuals": visual_count,
            "pictures": sum(1 for _ in _picture_shapes(slide)),
        })

        for phrase in banned_phrases:
            if phrase and phrase.casefold() in text.casefold():
                banned_count += 1
                issues.append(issue(_severity(config, "fail_on_banned_phrase"), "BANNED_PHRASE", f"Forbidden production-language phrase found: {phrase}", slide=slide_number))

        for shape in _walk_shapes(slide.shapes):
            for run in _iter_runs(shape):
                for typeface in _run_typefaces(run):
                    if allowed_fonts and typeface.casefold() not in allowed_fonts:
                        disallowed_font_count += 1
                        issues.append(issue(_severity(config, "fail_on_explicit_disallowed_font"), "DISALLOWED_FONT", f"Explicit font is outside allowed list: {typeface}", slide=slide_number, shape=_shape_name(shape), details={"font": typeface}))

        if slide_number not in bounds_exempt:
            for shape in slide.shapes:
                left, top, right, bottom = _rect(shape)
                if left < -bounds_tolerance or top < -bounds_tolerance or right > slide_width + bounds_tolerance or bottom > slide_height + bounds_tolerance:
                    bounds_count += 1
                    issues.append(issue(_severity(config, "fail_on_out_of_bounds"), "OUT_OF_BOUNDS", "Shape extends outside the slide canvas.", slide=slide_number, shape=_shape_name(shape), details={"left": left, "top": top, "right": right, "bottom": bottom, "slide_width": slide_width, "slide_height": slide_height}))

        if slide_number not in overlap_exempt:
            text_shapes = _top_level_text_shapes(slide)
            for index, first in enumerate(text_shapes):
                for second in text_shapes[index + 1 :]:
                    fraction = _overlap_fraction(first, second)
                    if fraction >= overlap_threshold:
                        overlap_count += 1
                        issues.append(issue(_severity(config, "fail_on_text_overlap"), "TEXT_GEOMETRY_OVERLAP", f"Text-bearing shapes overlap by {fraction:.1%} of the smaller box.", slide=slide_number, shape=f"{_shape_name(first)} <> {_shape_name(second)}", details={"fraction": fraction}))

        for picture in _picture_shapes(slide):
            try:
                digest = sha256_bytes(picture.image.blob)
                media_usage[digest].append({"slide": slide_number, "shape": _shape_name(picture), "filename": picture.image.filename, "px_size": list(picture.image.size)})
            except Exception as exc:
                issues.append(issue("warning", "PICTURE_HASH_ERROR", f"Picture bytes could not be read: {exc}", slide=slide_number, shape=_shape_name(picture)))
            if slide_number not in crop_exempt:
                crop_values = [float(picture.crop_left), float(picture.crop_right), float(picture.crop_top), float(picture.crop_bottom)]
                if any(abs(value) > crop_tolerance for value in crop_values):
                    crop_count += 1
                    issues.append(issue(_severity(config, "fail_on_nonzero_picture_crop"), "NONZERO_PICTURE_CROP", "Picture has non-zero PowerPoint crop values; complete-display evidence is not established.", slide=slide_number, shape=_shape_name(picture), details={"crop_left": crop_values[0], "crop_right": crop_values[1], "crop_top": crop_values[2], "crop_bottom": crop_values[3]}))

        if slide_number not in low_density_exempt:
            minimum = int(qa.get("minimum_body_characters_with_visual", 80) if visual_count else qa.get("minimum_body_characters_without_visual", 180))
            minimum_units = int(qa.get("minimum_substantive_units", 4))
            if body_characters < minimum or substantive_units < minimum_units:
                low_density_count += 1
                issues.append(issue(_severity(config, "fail_on_low_density"), "LOW_INFORMATION_DENSITY", f"Slide has {body_characters} characters, {substantive_units} substantive units and {visual_count} visuals; threshold is {minimum} characters and {minimum_units} units.", slide=slide_number))

    template_fraction = float(qa.get("ignore_template_media_repeated_on_fraction", 0.8))
    exempt_hashes = {value.casefold() for value in qa.get("duplicate_media_exempt_sha256", [])}
    template_hashes: set[str] = set()
    duplicate_hashes: dict[str, list[dict[str, Any]]] = {}
    for digest, usages in media_usage.items():
        distinct_slides = {usage["slide"] for usage in usages}
        if slide_count and len(distinct_slides) / slide_count >= template_fraction:
            template_hashes.add(digest)
        elif len(usages) > 1 and digest.casefold() not in exempt_hashes:
            duplicate_hashes[digest] = usages
            issues.append(issue(_severity(config, "fail_on_duplicate_evidence_media"), "DUPLICATE_EVIDENCE_MEDIA", f"The same non-template picture bytes are used {len(usages)} times on slides {sorted(distinct_slides)}.", details={"sha256": digest, "usages": usages}))

    page_config = presentation_config.get("page_number", {})
    page_number_results: list[dict[str, Any]] = []
    if page_config.get("required", False):
        start_slide = int(page_config.get("start_slide", 2))
        tolerance = int(Inches(float(page_config.get("center_tolerance_inches", 0.08))))
        bottom_fraction = float(page_config.get("bottom_zone_fraction", 0.35))
        for slide_number, slide in enumerate(presentation.slides, start=1):
            if slide_number < start_slide:
                continue
            candidates = []
            for shape in slide.shapes:
                page_text = _shape_text(shape).strip()
                if not re.fullmatch(r"\d+", page_text) or int(page_text) != slide_number:
                    continue
                center_x = int(shape.left) + int(shape.width) // 2
                center_y = int(shape.top) + int(shape.height) // 2
                if center_y < slide_height * (1 - bottom_fraction):
                    continue
                candidates.append((abs(center_x - slide_width // 2), center_x, _shape_name(shape)))
            if not candidates:
                issues.append(issue("error", "PAGE_NUMBER_MISSING", f"No bottom-zone page number text equal to {slide_number} was found.", slide=slide_number))
                page_number_results.append({"slide": slide_number, "found": False})
                continue
            distance, center_x, name = min(candidates)
            centered = distance <= tolerance
            page_number_results.append({"slide": slide_number, "found": True, "centered": centered, "distance_emu": distance, "shape": name})
            if not centered:
                issues.append(issue("error", "PAGE_NUMBER_NOT_CENTERED", f"Page number center is {distance / 914400:.3f} in from the slide width center; tolerance is {tolerance / 914400:.3f} in.", slide=slide_number, shape=name))

    asset_root = asset_root.resolve() if asset_root else None
    asset_hashes, metadata_hashes = _build_asset_hash_index(asset_root)
    content_hashes = set(media_usage) - template_hashes
    mapped_content_hashes = content_hashes & set(asset_hashes)
    metadata_mapped_hashes = content_hashes & set(metadata_hashes)
    unmapped_hashes = sorted(content_hashes - set(asset_hashes))
    if asset_root and not asset_root.is_dir():
        issues.append(issue("error", "ASSET_ROOT_NOT_FOUND", f"Asset root not found: {asset_root}"))
    elif asset_root:
        for digest in unmapped_hashes:
            issues.append(issue(_severity(config, "fail_on_unmapped_evidence_media"), "UNMAPPED_CONTENT_MEDIA", "A non-template picture embedded in the PPTX has no byte-identical file under the asset root.", details={"sha256": digest, "usages": media_usage[digest]}))

    metrics = {
        "pptx_sha256": sha256_file(pptx_path),
        "slide_count": slide_count,
        "slide_size_emu": [slide_width, slide_height],
        "slide_size_inches": [slide_width / 914400, slide_height / 914400],
        "crop_violation_count": crop_count,
        "out_of_bounds_count": bounds_count,
        "text_overlap_count": overlap_count,
        "banned_phrase_count": banned_count,
        "disallowed_font_count": disallowed_font_count,
        "low_density_slide_count": low_density_count,
        "embedded_picture_hash_count": len(media_usage),
        "template_picture_hash_count": len(template_hashes),
        "duplicate_non_template_hash_count": len(duplicate_hashes),
        "asset_mapping": {
            "asset_root": str(asset_root) if asset_root else None,
            "asset_file_hash_count": len(asset_hashes),
            "content_picture_hash_count": len(content_hashes),
            "mapped_content_picture_hash_count": len(mapped_content_hashes),
            "metadata_mapped_content_picture_hash_count": len(metadata_mapped_hashes),
            "unmapped_content_picture_hash_count": len(unmapped_hashes),
        },
        "page_numbers": page_number_results,
        "slides": slide_metrics,
        "issue_codes": dict(Counter(item["code"] for item in issues)),
    }
    return build_report("validate_pptx", pptx_path, issues, metrics)


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Audit a PPTX for structural, typography, evidence-media and density risks.")
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--config", type=Path)
    parser.add_argument("--asset-root", type=Path)
    parser.add_argument("--report", type=Path)
    args = parser.parse_args(argv)
    report = validate_pptx(args.pptx, args.config, args.asset_root)
    write_report(report, args.report)
    return 1 if report["status"] == "fail" else 0


if __name__ == "__main__":
    raise SystemExit(main())
