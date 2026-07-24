"""Build a self-contained five-slide academic presentation example.

The deck intentionally uses a blank canvas and original native PowerPoint
elements only. It contains no institutional template, private information or
third-party research figure.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).with_name("five_slide_academic_demo.pptx")
WIDE = 13.333333
HIGH = 7.5
FONT = "Microsoft YaHei"

NAVY = "102A43"
NAVY_DARK = "081A2C"
BLUE = "2E5E88"
GREEN = "2C7A5A"
GOLD = "C49A4A"
RED = "C63D3D"
TEXT = "172033"
MUTED = "5C6B7A"
PALE = "F3F6F8"
PALE_BLUE = "EAF1F7"
WHITE = "FFFFFF"
LINE = "CAD4DE"


def color(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def fill(shape, value: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(value)
    shape.line.fill.background()


def line(shape, value: str = LINE, width: float = 0.75) -> None:
    shape.line.color.rgb = color(value)
    shape.line.width = Pt(width)


def box(slide, x, y, w, h, value: str, line_color: str | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(shape, value)
    if line_color:
        line(shape, line_color)
    return shape


def text(
    slide,
    value: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    value_color: str = TEXT,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.0,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    paragraph.space_before = Pt(0)
    run = paragraph.add_run()
    run.text = value
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color(value_color)
    return shape


def section_header(slide, section: str, title: str, conclusion: str) -> None:
    box(slide, 0, 0, WIDE, 0.12, NAVY)
    text(slide, section, 0.65, 0.27, 1.8, 0.22, size=10.5, value_color=GREEN, bold=True)
    text(slide, title, 0.65, 0.52, 11.7, 0.42, size=27, value_color=NAVY, bold=True)
    box(slide, 0.65, 1.08, 0.72, 0.045, GOLD)
    text(slide, conclusion, 1.52, 1.01, 11.0, 0.27, size=13, value_color=MUTED)


def footer(slide, page: int, source: str) -> None:
    box(slide, 0.65, 6.77, 12.03, 0.012, LINE)
    text(slide, source, 0.65, 6.91, 5.75, 0.18, size=8.7, value_color=MUTED)
    text(slide, f"{page:02d}", 6.47, 6.91, 0.40, 0.20, size=9.5, value_color=NAVY, bold=True, align=PP_ALIGN.CENTER)


def add_cover(presentation: Presentation) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box(slide, 0, 0, WIDE, HIGH, NAVY_DARK)
    box(slide, 0.66, 1.08, 1.10, 0.055, GOLD)
    text(slide, "EVIDENCE-FIRST PPT WORKFLOW", 0.66, 1.35, 5.5, 0.25, size=12, value_color=GREEN, bold=True)
    text(slide, "学术PPT制作\n五页示例", 0.66, 1.86, 7.0, 1.25, size=39, value_color=WHITE, bold=True)
    text(slide, "从页面任务、素材证据到可交付文件的最小闭环", 0.70, 3.40, 8.0, 0.36, size=17, value_color="D9E4EC")
    for x, label, tone in [(0.70, "可复用", GREEN), (2.25, "可审计", GOLD), (3.80, "无机构模板", BLUE)]:
        box(slide, x, 4.28, 1.25 if label != "无机构模板" else 1.72, 0.36, tone)
        text(slide, label, x, 4.35, 1.25 if label != "无机构模板" else 1.72, 0.16, size=9.5, value_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(slide, "本示例只使用原创文字、原生表格与关系图；不含机构标识、个人信息、真实项目材料或第三方图件。", 0.70, 5.43, 10.8, 0.30, size=12, value_color="B8C7D6")
    box(slide, 0.70, 6.30, 5.5, 0.016, "365C7D")
    text(slide, "Source: Evidence-First PPT Workflow v1.0.2", 0.70, 6.55, 5.5, 0.18, size=8.7, value_color="9BB0C4")
    text(slide, "01", 6.47, 6.55, 0.40, 0.20, size=9.5, value_color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def add_agenda(presentation: Presentation) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    section_header(slide, "汇报提纲", "从制作顺序开始控制质量", "先冻结论证任务与证据边界，后决定素材和排版。")
    agenda = [
        ("一", "确定单页任务", "每页先写结论、证据与适用边界", RED),
        ("二", "审计真实素材", "确认图号、变量、版本、权利与完整性", NAVY),
        ("三", "构建图文关系", "让主视觉、文字解释和来源承担不同角色", NAVY),
        ("四", "执行QA与交付", "结构检查、全量渲染、人工审阅和哈希一致", NAVY),
    ]
    for index, (numeral, heading, body, tone) in enumerate(agenda):
        y = 1.62 + index * 0.83
        box(slide, 0.72, y + 0.04, 0.06, 0.55, tone)
        text(slide, numeral, 0.95, y, 0.35, 0.28, size=17, value_color=tone, bold=True)
        text(slide, heading, 1.42, y, 2.2, 0.25, size=16, value_color=tone, bold=True)
        text(slide, body, 1.42, y + 0.31, 4.95, 0.23, size=10.8, value_color=MUTED)
    box(slide, 7.08, 1.60, 0.03, 3.55, GOLD)
    text(slide, "本示例的制作约束", 7.42, 1.58, 4.8, 0.30, size=18, value_color=NAVY, bold=True)
    constraints = [
        "不把项目申请语言、内部制作提示或无来源图片写入页面；",
        "不以浅色卡片、放大标题或重复图形伪造“简洁”；",
        "所有正文使用统一字体；页码按页面几何中心定位；",
        "示例图形只解释PPT生产关系，不冒充真实科学观测。",
    ]
    for index, item in enumerate(constraints):
        y = 2.15 + index * 0.68
        text(slide, "•", 7.42, y, 0.22, 0.22, size=16, value_color=GREEN, bold=True)
        text(slide, item, 7.72, y, 4.55, 0.43, size=12.2, value_color=TEXT)
    text(slide, "导航规则：当前部分用红色高亮；进入下一部分前，目录再次出现并更新高亮。", 0.72, 5.65, 11.15, 0.30, size=12.5, value_color=NAVY, bold=True)
    footer(slide, 2, "Source: Evidence-First PPT Workflow v1.0.2 · 示例内容为原创制作规范。")


def add_page_spec(presentation: Presentation) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    section_header(slide, "一、页面规格", "先定义一页要完成的推理，再开始排版", "一页不是知识点堆栈；它应以一个可验证结论组织证据、解释与边界。")
    stages = [
        ("01", "受众问题", "观众此刻需要理解、比较或判断什么？", GREEN),
        ("02", "最小结论", "本页在不夸大事实的前提下能够证明什么？", GOLD),
        ("03", "证据资产", "哪一幅完整图、公式、表格或文本直接支撑结论？", BLUE),
        ("04", "适用边界", "参数、尺度、误差与引用应怎样限制结论？", RED),
    ]
    for index, (number, heading, body, tone) in enumerate(stages):
        y = 1.60 + index * 0.86
        box(slide, 0.72, y, 0.86, 0.63, tone)
        text(slide, number, 0.72, y + 0.19, 0.86, 0.20, size=12, value_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        text(slide, heading, 1.86, y + 0.02, 2.05, 0.24, size=15.5, value_color=NAVY, bold=True)
        text(slide, body, 1.86, y + 0.31, 3.35, 0.28, size=10.8, value_color=MUTED)
        if index < len(stages) - 1:
            box(slide, 1.12, y + 0.67, 0.05, 0.17, LINE)
    box(slide, 5.78, 1.60, 0.03, 3.67, GOLD)
    text(slide, "页面规格至少写清", 6.10, 1.56, 5.7, 0.26, size=18, value_color=NAVY, bold=True)
    items = [
        ("页面功能", "建立、比较、推导、诊断、验证或迁移"),
        ("核心内容", "不少于四条有实质含义的知识点或观察"),
        ("视觉中心", "一项主要证据，而不是多个互不相关的装饰"),
        ("页内引用", "图题、DOI/官方页面或方法文献可追溯"),
        ("验收条件", "完整显示、可读、无重叠、结论与证据匹配"),
    ]
    for index, (label, body) in enumerate(items):
        y = 2.08 + index * 0.58
        text(slide, label, 6.10, y, 1.20, 0.21, size=11.5, value_color=GREEN, bold=True)
        text(slide, body, 7.38, y, 4.65, 0.27, size=11.5, value_color=TEXT)
    text(slide, "结论：先冻结页面规格，能显著减少后期因为“图不对文、内容太空或标题越界”产生的返工。", 0.72, 5.70, 11.25, 0.30, size=12.5, value_color=NAVY, bold=True)
    footer(slide, 3, "Source: Evidence-First PPT Workflow v1.0.2 · 页面规格与内容密度标准。")


def add_asset_audit(presentation: Presentation) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    section_header(slide, "二、素材审计", "素材“相关”不足以成为页面证据", "每一幅进入PPT的期刊图或官方图，都要同时通过科学归属、显示完整性和使用权利审计。")
    table_x, table_y, table_w = 0.72, 1.62, 6.10
    columns = [(0.0, 1.42), (1.42, 2.05), (3.47, 2.63)]
    row_h = 0.58
    for row in range(5):
        y = table_y + row * row_h
        box(slide, table_x, y, table_w, row_h, NAVY if row == 0 else (PALE if row % 2 else WHITE), LINE)
        for start, width in columns:
            if start:
                box(slide, table_x + start, y, 0.010, row_h, LINE)
    headers = ["审计维度", "必须回答", "不通过时的动作"]
    for (start, width), label in zip(columns, headers):
        text(slide, label, table_x + start + 0.12, table_y + 0.16, width - 0.20, 0.20, size=11.0, value_color=WHITE, bold=True)
    rows = [
        ("归属与变量", "项目、图号、变量、单位是否明确？", "回到论文/官方页；不能猜测。"),
        ("时间与空间", "时间、区域、分辨率和产品性质是否明确？", "补充图题或更换素材。"),
        ("显示完整性", "坐标轴、图例、面板与必要标注是否齐全？", "找单图或完整图，不裁面板。"),
        ("访问与权利", "OA、订阅、仓储或许可依据是否已记录？", "仅保留元数据或请求许可。"),
    ]
    for row_index, row in enumerate(rows, start=1):
        y = table_y + row_index * row_h + 0.12
        for (start, width), label in zip(columns, row):
            text(slide, label, table_x + start + 0.12, y, width - 0.20, 0.30, size=10.1, value_color=TEXT, bold=(start == 0))
    box(slide, 7.30, 1.62, 0.03, 3.48, GREEN)
    text(slide, "论文与平台访问的优先顺序", 7.62, 1.58, 4.95, 0.27, size=18, value_color=NAVY, bold=True)
    routes = [
        ("1", "DOI与出版社正式页", "核对题名、版本、图号与正式图题"),
        ("2", "OA/官方原图/补充材料", "优先出版商提供的完整单图"),
        ("3", "机构订阅或作者仓储", "记录访问依据与内部使用边界"),
        ("4", "无法合法取得", "登记失败；不用缩略图或二次转载替代"),
    ]
    for index, (number, heading, body) in enumerate(routes):
        y = 2.05 + index * 0.70
        text(slide, number, 7.62, y, 0.26, 0.22, size=13, value_color=RED, bold=True)
        text(slide, heading, 8.00, y, 3.95, 0.21, size=12.3, value_color=NAVY, bold=True)
        text(slide, body, 8.00, y + 0.25, 4.28, 0.26, size=10.2, value_color=MUTED)
    text(slide, "访问控制不是技术障碍的“绕过目标”，而是素材权利与公开边界的一部分。", 0.72, 5.72, 11.15, 0.28, size=12.5, value_color=NAVY, bold=True)
    footer(slide, 4, "Source: Evidence-First PPT Workflow v1.0.2 · 真实素材与期刊访问规范。")


def add_qa(presentation: Presentation) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    section_header(slide, "三、构建与交付", "自动检查与人工放映共同决定“可交付”", "结构正确不等于放映可读；渲染好看也不等于来源、公式和版本已通过审计。")
    gates = [
        ("01", "PPTX结构", "页数、越界、文字重叠、字体、页码与裁剪值", GREEN),
        ("02", "素材证据", "文件哈希、图号、来源、变量、权利与使用范围", GOLD),
        ("03", "全量渲染", "PowerPoint导出、缩略图总览、重点页原尺寸检查", BLUE),
        ("04", "交付一致", "批准文件、交付副本、哈希、例外与已知限制", RED),
    ]
    for index, (number, heading, body, tone) in enumerate(gates):
        y = 1.62 + index * 0.83
        box(slide, 0.72, y, 0.76, 0.56, tone)
        text(slide, number, 0.72, y + 0.18, 0.76, 0.18, size=11.3, value_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        text(slide, heading, 1.76, y + 0.01, 1.65, 0.22, size=14, value_color=NAVY, bold=True)
        text(slide, body, 3.43, y + 0.02, 3.24, 0.30, size=10.6, value_color=TEXT)
        box(slide, 6.90, y + 0.19, 0.24, 0.18, tone)
        text(slide, "通过后进入下一道门禁", 7.34, y + 0.16, 3.35, 0.20, size=10.3, value_color=MUTED)
    box(slide, 10.70, 1.62, 0.03, 3.43, GOLD)
    text(slide, "不得交付的情形", 11.00, 1.59, 1.75, 0.28, size=16, value_color=NAVY, bold=True)
    risks = [
        "科研图被裁断或混入论文正文；",
        "公式以普通文本或错误符号模拟；",
        "页面信息不足、来源缺失或结论无证据；",
        "正式PPT与交付副本不是同一版本。",
    ]
    for index, item in enumerate(risks):
        y = 2.08 + index * 0.62
        text(slide, "×", 11.00, y, 0.22, 0.22, size=15, value_color=RED, bold=True)
        text(slide, item, 11.28, y, 1.25, 0.46, size=10.0, value_color=TEXT)
    text(slide, "交付定义：事实可追溯、图件完整、公式正确、页面可读、版本一致，并已完成五至八轮复核记录。", 0.72, 5.70, 11.18, 0.30, size=12.5, value_color=NAVY, bold=True)
    footer(slide, 5, "Source: Evidence-First PPT Workflow v1.0.2 · QA与交付验收标准。")


def build_deck(output: Path = OUT) -> Path:
    presentation = Presentation()
    presentation.slide_width = Inches(WIDE)
    presentation.slide_height = Inches(HIGH)
    core = presentation.core_properties
    core.title = "Evidence-First Academic PPT Five-Slide Demo"
    core.subject = "Synthetic, privacy-safe academic presentation workflow example"
    core.author = "Evidence-First PPT Workflow contributors"
    core.last_modified_by = "Evidence-First PPT Workflow contributors"
    core.keywords = "academic presentation; evidence-first; privacy-safe; example"
    core.comments = "No institutional template, private information, third-party figure, or user material is embedded."
    add_cover(presentation)
    add_agenda(presentation)
    add_page_spec(presentation)
    add_asset_audit(presentation)
    add_qa(presentation)
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)
    return output


if __name__ == "__main__":
    print(build_deck().resolve())
